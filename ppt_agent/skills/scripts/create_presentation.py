"""PowerPoint creation script - executed outside agent context.

This script contains the actual implementation logic for creating
PowerPoint presentations. It is called by the tool but its code
never enters the agent's context window (progressive disclosure).
"""

import os
import sys
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.oxml.xmlchemy import OxmlElement


def extract_numbers_from_text(text):
    """Extract numbers from text for potential chart data.

    Returns list of (label, value) tuples found in the text.
    """
    chart_data = []

    # Pattern 1: "Country: 123 medals" or "USA: 40 gold"
    # Captures: USA: 126 total medals, China: 91 total medals, etc.
    pattern1 = r'([A-Z][A-Za-z\s&]+):\s*(\d+)\s*(?:total\s+)?medals?'
    matches = re.findall(pattern1, text, re.IGNORECASE)
    for label, value in matches:
        label_clean = label.strip()
        if len(label_clean) > 1 and len(label_clean) < 30:  # Filter reasonable country names
            chart_data.append((label_clean, int(value)))

    # Pattern 2: More specific medal patterns
    # "40 gold, 44 silver, 42 bronze" -> Extract just the totals shown separately
    pattern2 = r'([A-Z][A-Za-z\s]+).*?(\d+)\s+gold.*?(\d+)\s+silver.*?(\d+)\s+bronze'
    matches = re.findall(pattern2, text, re.IGNORECASE)
    for match in matches:
        if len(match) == 4:
            country, gold, silver, bronze = match
            total = int(gold) + int(silver) + int(bronze)
            if country.strip() and total > 0:
                chart_data.append((country.strip(), total))

    # Pattern 3: "10,500 athletes" or "4,448 GW" type stats
    pattern3 = r'(\d+(?:[,\.]\d+)?)\s*([A-Z]{1,3})?[\s:]*(athletes?|countries?|nations?|events?|sports?|participants?|capacity|billion|million|thousand|GW|MW|TW|vehicles?|sales?)'
    matches = re.findall(pattern3, text, re.IGNORECASE)
    for match in matches[:10]:
        if len(match) == 3:
            value_str, unit, label = match
            try:
                clean_value = float(value_str.replace(',', ''))
                # Adjust for units
                if unit and unit.upper() in ['TW']:
                    clean_value = int(clean_value * 1000)  # TW to GW
                    label = f"{label} (GW)"
                elif clean_value >= 1 and clean_value < 1000000:
                    chart_data.append((label.capitalize(), int(clean_value)))
            except:
                pass

    # Pattern 4: "2022: $384.65 billion" or "China: 45%"
    pattern4 = r'([A-Z][A-Za-z\s]+|20\d{2}):\s*(?:\$|€)?(\d+(?:[,\.]\d+)?)\s*(billion|million|%)?'
    matches = re.findall(pattern4, text, re.IGNORECASE)
    for label, value, unit in matches[:10]:
        try:
            clean_value = float(value.replace(',', ''))
            if unit and 'billion' in unit.lower():
                label = f"{label.strip()} ($B)"
            elif unit and '%' in unit:
                label = f"{label.strip()} (%)"
            if clean_value > 0 and len(label.strip()) > 1:
                chart_data.append((label.strip(), int(clean_value)))
        except:
            pass

    # Deduplicate by label (keep first occurrence)
    seen = set()
    unique_data = []
    for label, value in chart_data:
        if label not in seen and value > 0:
            seen.add(label)
            unique_data.append((label, value))

    return unique_data[:8]  # Limit to top 8 for chart clarity


def create_bar_chart_slide(prs, title, chart_data_list, bullet_slide_layout):
    """Create a slide with a bar chart.

    Args:
        prs: Presentation object
        title: Slide title
        chart_data_list: List of (label, value) tuples
        bullet_slide_layout: Layout to use
    """
    if not chart_data_list or len(chart_data_list) < 2:
        return None  # Need at least 2 data points for a chart

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    # Set title
    title_shape = shapes.title
    title_shape.text = title

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in chart_data_list]
    chart_data.add_series('Count', [item[1] for item in chart_data_list])

    # Add chart
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Configure chart
    chart.has_legend = False
    chart.chart_title.text_frame.text = title

    return slide


def create_pie_chart_slide(prs, title, chart_data_list, bullet_slide_layout):
    """Create a slide with a pie chart.

    Args:
        prs: Presentation object
        title: Slide title
        chart_data_list: List of (label, value) tuples
        bullet_slide_layout: Layout to use
    """
    if not chart_data_list or len(chart_data_list) < 2:
        return None

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    # Set title
    title_shape = shapes.title
    title_shape.text = title

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in chart_data_list]
    chart_data.add_series('Distribution', [item[1] for item in chart_data_list])

    # Add pie chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    # Configure chart
    chart.has_legend = True
    chart.chart_title.text_frame.text = title

    # Enable data labels showing percentages
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False

    return slide


def create_table_slide(prs, title, table_data, bullet_slide_layout):
    """Create a slide with a table.

    Args:
        prs: Presentation object
        title: Slide title
        table_data: List of lists [[header1, header2], [row1col1, row1col2], ...]
        bullet_slide_layout: Layout to use
    """
    if not table_data or len(table_data) < 2:
        return None

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    # Set title
    title_shape = shapes.title
    title_shape.text = title

    # Calculate table dimensions
    rows = len(table_data)
    cols = len(table_data[0])

    # Add table
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    table = shapes.add_table(rows, cols, x, y, cx, cy).table

    # Populate table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_value)

            # Style header row
            if row_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(14)

    return slide


def parse_research_for_visuals(research_data):
    """Parse research data to extract visual suggestions and data.

    Returns dict with:
    - 'bar_charts': List of (title, data) for bar charts
    - 'pie_charts': List of (title, data) for pie charts
    - 'tables': List of (title, table_data) for tables
    - 'suggestions': List of suggested visualizations
    """
    visuals = {
        'bar_charts': [],
        'pie_charts': [],
        'tables': [],
        'suggestions': []
    }

    if not research_data:
        return visuals

    lines = research_data.split('\n')

    # Extract visual suggestions section
    in_visual_section = False
    for line in lines:
        line = line.strip()

        if '**Visual Suggestions:**' in line or '**Visualizations:**' in line:
            in_visual_section = True
            continue

        if in_visual_section:
            if line.startswith('**') and line.endswith('**'):
                in_visual_section = False
            elif line.startswith('-'):
                suggestion = line.lstrip('- ').strip()
                if suggestion:
                    visuals['suggestions'].append(suggestion)

    # Try to extract chart-worthy data
    chart_data = extract_numbers_from_text(research_data)

    if chart_data:
        # Create a bar chart for country comparisons (if we have multiple countries/entities)
        country_data = [(label, value) for label, value in chart_data
                       if not any(word in label.lower() for word in ['athletes', 'countries', 'events', 'sports', 'participants'])]

        if len(country_data) >= 2:
            visuals['bar_charts'].append(("Medal Count by Country", country_data[:8]))

        # Create a pie chart for medal type distribution if we can find gold/silver/bronze
        # Pattern: Look for "40 gold, 44 silver, 42 bronze" in text
        medal_types_match = re.search(r'(\d+)\s+gold.*?(\d+)\s+silver.*?(\d+)\s+bronze', research_data, re.IGNORECASE)
        if medal_types_match:
            gold, silver, bronze = medal_types_match.groups()
            medal_type_data = [
                ('Gold', int(gold)),
                ('Silver', int(silver)),
                ('Bronze', int(bronze))
            ]
            visuals['pie_charts'].append(("Medal Type Distribution", medal_type_data))

        # Add summary stats as a potential table (participation data)
        summary_stats = [(label, value) for label, value in chart_data
                        if any(word in label.lower() for word in ['athletes', 'countries', 'events', 'sports', 'participants'])]

        if len(summary_stats) >= 2:
            # Create table data
            table_data = [['Metric', 'Count']]  # Header
            for label, value in summary_stats[:5]:
                table_data.append([label, str(value)])
            visuals['tables'].append(("Event Statistics", table_data))

    return visuals


def create_powerpoint(
    topic: str,
    num_slides: int = 5,
    include_title_slide: bool = True,
    output_dir: str = "./output",
    research_data: str = None,
) -> dict:
    """Create a PowerPoint presentation with optional research-enriched content.

    This function is called by the tool wrapper but its implementation
    stays outside the agent's context.

    Args:
        topic: Main topic/subject of the presentation
        num_slides: Number of content slides (excluding title)
        include_title_slide: Whether to include a title slide
        output_dir: Directory to save the file
        research_data: Optional research findings to incorporate into slides

    Returns:
        Dictionary with success status, filepath, and message
    """
    try:
        # Validate inputs
        if num_slides < 1 or num_slides > 20:
            return {
                "success": False,
                "error": "num_slides must be between 1 and 20",
                "filepath": None,
            }

        if not topic or len(topic.strip()) == 0:
            return {
                "success": False,
                "error": "topic cannot be empty",
                "filepath": None,
            }

        # Create output directory
        os.makedirs(output_dir, exist_ok=True)

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Add title slide
        if include_title_slide:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = topic
            subtitle_text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            if research_data:
                subtitle_text += " | Research-Enhanced Presentation"
            subtitle.text = subtitle_text

        # Add content slides
        bullet_slide_layout = prs.slide_layouts[1]

        # If research data is provided, add a research findings slide first
        if research_data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = "Key Research Findings"

            body_shape = shapes.placeholders[1]
            text_frame = body_shape.text_frame

            # Parse research data into bullet points
            # Simple parsing: split by newlines and filter
            research_lines = [line.strip() for line in research_data.split('\n') if line.strip() and not line.strip().startswith('#')]

            if research_lines:
                text_frame.text = research_lines[0]
                for line in research_lines[1:8]:  # Limit to 8 bullets to avoid overcrowding
                    if line and line != '**' and not line.startswith('**'):
                        p = text_frame.add_paragraph()
                        # Clean up markdown-style formatting
                        clean_line = line.replace('**', '').replace('- ', '').replace('* ', '')
                        p.text = clean_line
                        p.level = 0 if line.startswith('-') or line.startswith('*') else 1
            else:
                text_frame.text = "Research findings incorporated throughout presentation"

        # Parse research for visual content (charts, tables)
        visuals = parse_research_for_visuals(research_data) if research_data else {
            'bar_charts': [], 'pie_charts': [], 'tables': [], 'suggestions': []
        }

        # Add bar chart slides
        for chart_title, chart_data in visuals.get('bar_charts', []):
            print(f"[DEBUG] Creating bar chart: {chart_title} with {len(chart_data)} data points", file=sys.stderr)
            create_bar_chart_slide(prs, chart_title, chart_data, bullet_slide_layout)

        # Add pie chart slides
        for chart_title, chart_data in visuals.get('pie_charts', []):
            print(f"[DEBUG] Creating pie chart: {chart_title} with {len(chart_data)} segments", file=sys.stderr)
            create_pie_chart_slide(prs, chart_title, chart_data, bullet_slide_layout)

        # Add table slides
        for table_title, table_data in visuals.get('tables', []):
            print(f"[DEBUG] Creating table: {table_title} with {len(table_data)} rows", file=sys.stderr)
            create_table_slide(prs, table_title, table_data, bullet_slide_layout)

        # Parse research data into usable bullets for content slides
        research_bullets = []

        # DEBUG: Log research_data state
        print(f"[DEBUG] research_data received: {research_data is not None}", file=sys.stderr)
        print(f"[DEBUG] research_data length: {len(research_data) if research_data else 0}", file=sys.stderr)
        if research_data:
            print(f"[DEBUG] research_data preview (first 200 chars): {research_data[:200]}", file=sys.stderr)

        if research_data:
            # Extract all meaningful bullet points from research data
            # Split by lines and filter for actual content
            lines = research_data.split('\n')

            for line in lines:
                line = line.strip()

                # Skip empty lines, headers, and URLs
                if not line or line.startswith('#') or line.startswith('http'):
                    continue

                # Skip section markers
                if line.startswith('**') or line.endswith('**'):
                    continue

                # Clean up markdown bullet markers
                cleaned = line
                for marker in ['- ', '* ', '• ']:
                    if cleaned.startswith(marker):
                        cleaned = cleaned[len(marker):]

                cleaned = cleaned.strip()

                # Only keep substantial content (not section headers, not URLs)
                if (len(cleaned) > 15 and
                    not cleaned.lower().startswith(('key findings', 'visual suggestions', 'sources', 'key insight')) and
                    not cleaned.startswith('http') and
                    ':' not in cleaned[:30]):  # Likely a data point, not a header
                    research_bullets.append(cleaned)

        # DEBUG: Log parsing results
        print(f"[DEBUG] research_bullets extracted: {len(research_bullets)} bullets", file=sys.stderr)
        if research_bullets:
            print(f"[DEBUG] First 3 bullets: {research_bullets[:3]}", file=sys.stderr)

        # Add regular content slides
        for i in range(1, num_slides + 1):
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            # Set title
            title_shape = shapes.title
            title_shape.text = f"{topic} - Point {i}"

            # Add content
            body_shape = shapes.placeholders[1]
            text_frame = body_shape.text_frame

            # DEBUG: Log slide content decision
            print(f"[DEBUG] Slide {i} - research_data exists: {research_data is not None}, research_bullets count: {len(research_bullets)}", file=sys.stderr)

            if research_data and research_bullets:
                # Distribute research bullets across slides
                # Calculate which bullets go on this slide (distribute evenly)
                bullets_per_slide = max(3, len(research_bullets) // num_slides)
                start_idx = (i - 1) * bullets_per_slide
                end_idx = min(start_idx + bullets_per_slide, len(research_bullets))

                slide_bullets = research_bullets[start_idx:end_idx]

                if slide_bullets:
                    # Add first bullet as main text
                    text_frame.text = slide_bullets[0]
                    # Add remaining bullets as paragraphs
                    for bullet in slide_bullets[1:]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                else:
                    # Fallback if no bullets for this slide
                    text_frame.text = f"Additional insights related to {topic}"
            else:
                # No research data - use generic content
                text_frame.text = f"Key concept {i} related to {topic}"
                # Add bullet points
                for j in range(3):
                    p = text_frame.add_paragraph()
                    p.text = f"Supporting detail {j + 1} for concept {i}"
                    p.level = 1

        # Generate filename
        safe_topic = "".join(c if c.isalnum() or c.isspace() else "_" for c in topic)
        safe_topic = safe_topic.replace(" ", "_")[:50]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_topic}_{timestamp}.pptx"
        filepath = os.path.join(output_dir, filename)

        # Save presentation
        prs.save(filepath)

        total_slides = num_slides + (1 if include_title_slide else 0)

        return {
            "success": True,
            "filepath": filepath,
            "total_slides": total_slides,
            "message": f"PowerPoint presentation created successfully: {filepath} (Total slides: {total_slides})",
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "filepath": None,
            "message": f"Error creating PowerPoint presentation: {str(e)}",
        }
